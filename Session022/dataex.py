import requests
from bs4 import BeautifulSoup
import re
from googlesearch import search
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def fetch_data(url):
    try:
        response = requests.get(url)
        response.raise_for_status() 
        return response.text
    except requests.RequestException as e:
        print(f"Error fetching data from {url}: {e}")
        return None

def extract_numerical_data(content, keywords):
    soup = BeautifulSoup(content, 'html.parser')
    text = soup.get_text()

    # Regular expression to find numbers (both integers and floats)
    numerical_data = re.findall(r'\b\d+(?:\.\d+)?\b', text)

    # Dictionary to hold keyword and corresponding numerical data
    keyword_data = {keyword: [] for keyword in keywords}

    # Filter based on keywords
    for keyword in keywords:
        if keyword in text.lower():
            # Add all found numbers to the corresponding keyword
            keyword_data[keyword].extend(numerical_data)

    return {k: list(set(v)) for k, v in keyword_data.items() if v}  # Return non-empty keyword data


def fetch_data(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        return response.text
    except requests.RequestException as e:
        print(f"Error fetching data from {url}: {e}")
        return None

def fetch_numerical_data(keyword):
    # Search for the query and get 15 results
    urls = search(keyword, num_results=15)
    numerical_data = []

    for url in urls:
        print(f"Fetching data from: {url}")
        content = fetch_data(url)
        if content:
            # Extract numerical data from content
            numbers = re.findall(r'\b\d+(?:\.\d+)?\b', content)
            numerical_data.extend(numbers)

    return list(set(numerical_data))  # Return unique numerical values

def extract_data(numerical_data):
    extracted_info = {
        "years": [],
        "sales": [],
        "market": [],
        "production": []
    }

    for value in numerical_data:
        str_value = str(value)  # Convert to string to handle all cases

        # Check for year (four-digit number)
        if str_value.isdigit() and len(str_value) == 4 and 2000 <= int(str_value) <= 2024:
            extracted_info["years"].append(str_value)

        # Check for sales (numbers typically greater than a certain threshold)
        elif str_value.replace('.', '', 1).isdigit() and float(str_value) >= 1000:  # Example threshold for sales
            extracted_info["sales"].append(str_value)

        # Check for market share (assume percentage representation)
        elif str_value.isdigit() and 0 <= int(str_value) <= 100:  # Example for market share percentage
            extracted_info["market"].append(str_value)

        # Check for production (assume larger numbers are production counts)
        elif str_value.isdigit() and int(str_value) > 1000:  # Example threshold for production
            extracted_info["production"].append(str_value)

    return extracted_info

def main():
    # Define a keyword for fetching numerical data
    keyword = input("Enter a product or query you want to search for: ")
    
    # Step 1: Fetch numerical data based on the keyword
    numerical_data = fetch_numerical_data(keyword)

    # Step 2: Extract Data
    extracted_info = extract_data(numerical_data)

    # Step 3: Convert extracted_info to DataFrame for further processing
    df_years = pd.DataFrame(extracted_info['years'], columns=['Years'])
    df_sales = pd.DataFrame(extracted_info['sales'], columns=['Sales'])
    df_market = pd.DataFrame(extracted_info['market'], columns=['Market'])
    df_production = pd.DataFrame(extracted_info['production'], columns=['Production'])

    # Combine all into a single DataFrame if needed
    combined_df = pd.concat([df_years, df_sales, df_market, df_production], axis=1)

    # Print the combined DataFrame
    print("Extracted Data:")
    print(combined_df)

'task 2'


def read_data(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == '.csv':
        df = pd.read_csv(file_path)
    elif file_extension in ['.xls', '.xlsx']:
        df = pd.read_excel(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")
    return df

def detect_field_types(df):
    numeric_columns = df.select_dtypes(include=['number']).columns.tolist()
    categorical_columns = df.select_dtypes(exclude=['number']).columns.tolist()
    return numeric_columns, categorical_columns

def generate_graphs(df):
    numeric_columns, categorical_columns = detect_field_types(df)
    
    # Generate histograms for numeric columns
    for col in numeric_columns:
        plt.figure(figsize=(8, 6))
        sns.histplot(df[col].dropna(), kde=True, color='blue')
        plt.title(f'Histogram of {col}', fontsize=16)
        plt.xlabel(col, fontsize=14)
        plt.ylabel('Frequency', fontsize=14)
        plt.grid(True)
        plt.tight_layout()
        plt.show()
    
    # Generate bar plots for categorical columns
    for col in categorical_columns:
        if df[col].nunique() < 20:
            plt.figure(figsize=(10, 6))
            sns.countplot(y=df[col], data=df, hue=df[col], palette='Set2')
            plt.title(f'Bar Plot of {col}', fontsize=16)
            plt.ylabel(col, fontsize=14)
            plt.xlabel('Count', fontsize=14)
            plt.grid(True)
            plt.tight_layout()
            plt.show()

def calculate_insights(df):
    insights = {}
    numeric_columns, categorical_columns = detect_field_types(df)

    # Calculate insights for numeric columns
    for col in numeric_columns:
        insights[col] = {
            'Mean': round(df[col].mean(), 2),
            'Median': round(df[col].median(), 2),
            'Standard Deviation': round(df[col].std(), 2),
            'Min': df[col].min(),
            'Max': df[col].max()
        }

    # Calculate insights for categorical columns
    for col in categorical_columns:
        insights[col] = {
            'Unique Values': df[col].nunique(),
            'Most Frequent Value': df[col].mode().iloc[0]
        }

    return insights

def add_insights_to_ppt(prs, insights):
    slide_layout = prs.slide_layouts[5]  # Title and Content slide
    for col, stats in insights.items():
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f'Insights for {col}'

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5)

        for stat, value in stats.items():
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f'{stat}: {value}'
            p.font.size = Pt(14)
            tf.alignment = PP_ALIGN.LEFT
            top += height

def save_to_ppt(df):
    numeric_columns, categorical_columns = detect_field_types(df)
    insights = calculate_insights(df)
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide

    # Save numeric graphs
    for col in numeric_columns:
        fig = plt.figure(figsize=(8, 6))
        sns.histplot(df[col].dropna(), kde=True, color='blue')
        plt.title(f'Histogram of {col}', fontsize=16)
        plt.xlabel(col, fontsize=14)
        plt.ylabel('Frequency', fontsize=14)
        plt.grid(True)
        plt.tight_layout()
        fig.savefig('temp.png', bbox_inches='tight', dpi=300)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture('temp.png', Inches(1), Inches(1), width=Inches(8))
        plt.close(fig)
        os.remove('temp.png')

    # Save categorical graphs
    for col in categorical_columns:
        if df[col].nunique() < 20:
            fig = plt.figure(figsize=(10, 6))
            sns.countplot(y=df[col], data=df, hue=df[col], palette='Set2')
            plt.title(f'Bar Plot of {col}', fontsize=16)
            plt.ylabel(col, fontsize=14)
            plt.xlabel('Count', fontsize=14)
            plt.grid(True)
            plt.tight_layout()
            fig.savefig('temp.png', bbox_inches='tight', dpi=300)
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture('temp.png', Inches(1), Inches(1), width=Inches(10))
            plt.close(fig)
            os.remove('temp.png')

    add_insights_to_ppt(prs, insights)
    prs.save('output4.pptx')

def main():
    query = input("Enter the product or query you want to search for: ").lower()

    # Search for the query and get 15 results
    urls = search(query, num_results=15)
    ##    print(urls)
    keywords = ['sales', 'production', 'market', 'revenue', 'profit', 'loss']
    all_keyword_data = {}

    for url in urls:
        print(f"Fetching data from: {url}")
        content = fetch_data(url)
        if content:
            keyword_data = extract_numerical_data(content, keywords)
            for key, values in keyword_data.items():
                if key not in all_keyword_data:
                    all_keyword_data[key] = []
                all_keyword_data[key].extend(values)

    # Print results
    if all_keyword_data:
        for keyword, values in all_keyword_data.items():
            print(f"Numerical data for '{keyword}': {list(set(values))}")
    else:
        print("No numerical data found.")
    file_path = input("Enter CSV or Excel file: ")
    df = read_data(file_path)
    generate_graphs(df)
    save_to_ppt(df)

if __name__ == '__main__':
    main()
