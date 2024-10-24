import requests
from bs4 import BeautifulSoup
import re
from googlesearch import search
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# Fetch data from URL with error handling
def fetch_data(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        return response.text
    except requests.RequestException as e:
        print(f"Error fetching data from {url}: {e}")
        return None


# Extract numerical data based on given keywords from content
def extract_numerical_data(content, keywords):
    soup = BeautifulSoup(content, "html.parser")
    text = soup.get_text()
    numerical_data = re.findall(r"\b\d+(?:\.\d+)?\b", text)
    keyword_data = {keyword: [] for keyword in keywords}

    for keyword in keywords:
        if keyword in text.lower():
            keyword_data[keyword].extend(numerical_data)

    return {k: list(set(v)) for k, v in keyword_data.items() if v}


# Main logic for searching and fetching numerical data from URLs
def search_and_extract_data(query, keywords):
    urls = search(query, num_results=15)
    all_keyword_data = {keyword: [] for keyword in keywords}

    for url in urls:
        content = fetch_data(url)
        if content:
            keyword_data = extract_numerical_data(content, keywords)
            for key, values in keyword_data.items():
                all_keyword_data[key].extend(values)

    return {k: list(set(v)) for k, v in all_keyword_data.items() if v}


# Function to extract and categorize numerical data
def extract_data(numerical_data):
    extracted_info = {"years": [], "sales": [], "market": [], "production": []}

    for value in numerical_data:
        str_value = str(value)

        if (
            str_value.isdigit()
            and len(str_value) == 4
            and 2000 <= int(str_value) <= 2024
        ):
            extracted_info["years"].append(str_value)
        elif str_value.replace(".", "", 1).isdigit() and float(str_value) >= 1000:
            extracted_info["sales"].append(str_value)
        elif str_value.isdigit() and 0 <= int(str_value) <= 100:
            extracted_info["market"].append(str_value)
        elif str_value.isdigit() and int(str_value) > 1000:
            extracted_info["production"].append(str_value)

    return extracted_info


# Read CSV or Excel file
def read_data(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".csv":
        return pd.read_csv(file_path)
    elif ext in [".xls", ".xlsx"]:
        return pd.read_excel(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


# Detect numeric and categorical columns
def detect_field_types(df):
    return (
        df.select_dtypes(include=["number"]).columns.tolist(),
        df.select_dtypes(exclude=["number"]).columns.tolist(),
    )


# Generate graphs based on numeric and categorical columns
def generate_graphs(df):
    numeric_columns, categorical_columns = detect_field_types(df)

    for col in numeric_columns:
        plt.figure(figsize=(8, 6))
        sns.histplot(df[col].dropna(), kde=True, color="blue")
        plt.title(f"Histogram of {col}", fontsize=16)
        plt.xlabel(col, fontsize=14)
        plt.ylabel("Frequency", fontsize=14)
        plt.grid(True)
        plt.tight_layout()
        plt.show()

    for col in categorical_columns:
        if df[col].nunique() < 20:
            plt.figure(figsize=(10, 6))
            sns.countplot(y=df[col], data=df, palette="Set2")
            plt.title(f"Bar Plot of {col}", fontsize=16)
            plt.ylabel(col, fontsize=14)
            plt.xlabel("Count", fontsize=14)
            plt.grid(True)
            plt.tight_layout()
            plt.show()


# Calculate statistics insights for numeric and categorical columns
def calculate_insights(df):
    insights = {}
    numeric_columns, categorical_columns = detect_field_types(df)

    for col in numeric_columns:
        insights[col] = {
            "Mean": round(df[col].mean(), 2),
            "Median": round(df[col].median(), 2),
            "Std Dev": round(df[col].std(), 2),
            "Min": df[col].min(),
            "Max": df[col].max(),
        }

    for col in categorical_columns:
        insights[col] = {
            "Unique Values": df[col].nunique(),
            "Most Frequent": df[col].mode().iloc[0],
        }

    return insights


# Add insights to PowerPoint slides
def add_insights_to_ppt(prs, insights):
    slide_layout = prs.slide_layouts[5]
    for col, stats in insights.items():
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Insights for {col}"
        left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(0.5)

        for stat, value in stats.items():
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            para = text_frame.add_paragraph()
            para.text = f"{stat}: {value}"
            para.font.size = Pt(14)
            text_frame.alignment = PP_ALIGN.LEFT
            top += height


# Save graphs and insights into a PowerPoint file
def save_to_ppt(df):
    numeric_columns, categorical_columns = detect_field_types(df)
    insights = calculate_insights(df)
    prs = Presentation()

    slide_layout = prs.slide_layouts[6]
    for col in numeric_columns:
        fig = plt.figure(figsize=(8, 6))
        sns.histplot(df[col].dropna(), kde=True, color="blue")
        plt.title(f"Histogram of {col}", fontsize=16)
        plt.xlabel(col, fontsize=14)
        plt.ylabel("Frequency", fontsize=14)
        plt.grid(True)
        plt.tight_layout()
        fig.savefig("temp.png", bbox_inches="tight", dpi=300)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture("temp.png", Inches(1), Inches(1), width=Inches(8))
        plt.close(fig)
        os.remove("temp.png")

    for col in categorical_columns:
        if df[col].nunique() < 20:
            fig = plt.figure(figsize=(10, 6))
            sns.countplot(y=df[col], data=df, palette="Set2")
            plt.title(f"Bar Plot of {col}", fontsize=16)
            plt.ylabel(col, fontsize=14)
            plt.xlabel("Count", fontsize=14)
            plt.grid(True)
            plt.tight_layout()
            fig.savefig("temp.png", bbox_inches="tight", dpi=300)
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture("temp.png", Inches(1), Inches(1), width=Inches(10))
            plt.close(fig)
            os.remove("temp.png")

    add_insights_to_ppt(prs, insights)
    prs.save("output_presentation.pptx")


# Main function
def main():
    file_path = input("Enter CSV or Excel file: ")
    df = read_data(file_path)
    generate_graphs(df)
    save_to_ppt(df)


if __name__ == "__main__":
    main()
